"""V3 Part 1: Cover + Overview (6 slides) - 炫酷深蓝科技风 + 详尽备注"""
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN as A
from pptx.enum.shapes import MSO_SHAPE as S
# 金融科技配色：深蓝底+青色主色+金色强调
BG1=C(0x06,0x12,0x22);BG2=C(0x0A,0x1E,0x35);CY=C(0,0xD4,0xFF);GO=C(0xFF,0xB8,0);GR=C(0,0xE6,0x76);PU=C(0xA8,0x5C,0xFF);W=C(255,255,255);G=C(0x6B,0x7B,0x8D);LG=C(0xA0,0xB0,0xC0);F='Calibri'
def bg(s,c=BG1):f=s.background.fill;f.solid();f.fore_color.rgb=c
def t(s,l,top,w,h,tx,sz=14,c=W,b=False,a=A.LEFT):bx=s.shapes.add_textbox(l,top,w,h);tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=tx;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.font.name=F;p.alignment=a;return bx
def mt(s,l,top,w,h,lines):
 bx=s.shapes.add_textbox(l,top,w,h);tf=bx.text_frame;tf.word_wrap=True
 for i,ld in enumerate(lines):
  p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.text=ld[0];p.font.size=Pt(ld[1] if len(ld)>1 else 14);p.font.color.rgb=ld[2] if len(ld)>2 else W;p.font.bold=ld[3] if len(ld)>3 else False;p.font.name=F;p.alignment=ld[4] if len(ld)>4 else A.LEFT;p.space_after=Pt(3)
 return bx
def bar(s,l,top,w,h=Pt(3),c=CY):sh=s.shapes.add_shape(S.RECTANGLE,l,top,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
def card(s,l,top,w,h,c=BG2):sh=s.shapes.add_shape(S.ROUNDED_RECTANGLE,l,top,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
def hdr(s,title,sub=""):
 t(s,I(.8),I(.3),I(14),I(.6),title,30,W,True);bar(s,I(.8),I(.9),I(14.4),Pt(3),CY)
 if sub:t(s,I(.8),I(1.05),I(14),I(.35),sub,12,G)
def ftr(s):t(s,I(.8),I(8.55),I(14),I(.3),"华锐技术 Archforce Technology | 深度技术研究 V3 | 2026.03",8,G)
def n(s,tx):s.notes_slide.notes_text_frame.text=tx
def metric(s,l,top,num,label,c=CY):
 card(s,l,top,I(2.3),I(1.1),BG2)
 t(s,l+I(.1),top+I(.1),I(2.1),I(.55),str(num),34,c,True,A.CENTER)
 t(s,l+I(.1),top+I(.65),I(2.1),I(.3),label,10,LG,False,A.CENTER)
prs=Presentation();prs.slide_width=I(16);prs.slide_height=I(9)

# S1 Cover - 炫酷封面
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,BG1)
bar(s,I(0),I(0),I(16),Pt(4),CY)
bar(s,I(0),I(8.95),I(16),Pt(4),CY)
t(s,I(1),I(1.8),I(14),I(.5),"DEEP DIVE RESEARCH",18,CY,True,A.CENTER)
t(s,I(1),I(2.5),I(14),I(1.2),"华锐技术",52,W,True,A.CENTER)
t(s,I(1),I(3.7),I(14),I(.8),"Archforce Technology",36,CY,False,A.CENTER)
bar(s,I(5),I(4.7),I(6),Pt(3),GO)
t(s,I(1),I(5.2),I(14),I(.5),"中国领先的金融科技公司 · 分布式低时延技术引领者",16,LG,False,A.CENTER)
t(s,I(1),I(6),I(14),I(.4),"60ns消息传输 · 110+券商 · TOP20覆盖95% · 百亿量化私募100%",13,GO,False,A.CENTER)
t(s,I(1),I(7),I(14),I(.4),"2026年3月 | 内部研究材料",13,G,False,A.CENTER)
ftr(s)
n(s,"""【演讲开场白】
各位好，今天我来给大家深入介绍一家非常重要的金融科技公司——华锐技术，英文名Archforce Technology。

为什么说它重要？先给大家几个数字感受一下：
- 消息中间件时延60纳秒——这是什么概念？CPU访问L3缓存大约40纳秒，华锐几乎做到了零开销传输
- 交易平台覆盖了110多家券商，TOP20券商覆盖率95%
- 百亿级量化私募覆盖率100%——量化交易对时延最敏感，全部选择华锐

今天的演讲将从公司概况、技术架构、产品体系、市场地位等多个维度，带大家深入了解华锐的全貌。这份报告基于官网全站30个页面的深度抓取，以及CB Insights、PitchBook、Forbes China等第三方数据源的交叉验证。

接下来我们先看目录，了解今天要讲的内容框架。""")

# S2 Agenda
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
hdr(s,"目录 AGENDA")
items=[("01","公司概况与融资","Company Overview",CY),("02","基础软件全栈","AMI · AMQ · AMDB · Platform · Euler",CY),("03","核心交易平台 ATP","T1~T9 Full Product Line",GO),("04","实时风控 ARC","30μs / 200 Rules",GO),("05","高速行情 AMD","FPGA 130ns Decoding",GR),("06","业务平台群","ACP · AIM · ATA · ATH",GR),("07","信创 · 国际化 · 市场地位","Xinchuang & Globalization",PU),("08","总结与展望","Summary & Outlook",PU)]
for i,(num,zh,en,c) in enumerate(items):
 y=I(1.4+i*.85);card(s,I(.8),y,I(14.4),I(.75))
 t(s,I(1.1),y+I(.12),I(.8),I(.5),num,22,c,True)
 bar(s,I(2),y+I(.15),Pt(3),I(.45),c)
 t(s,I(2.3),y+I(.12),I(6),I(.5),zh,17,W,True)
 t(s,I(9),y+I(.17),I(6),I(.4),en,12,G)
ftr(s)
n(s,"""【目录介绍】
今天的演讲分为8个章节，我简单介绍一下每章要讲什么：

第1章 公司概况与融资——了解华锐是谁、做什么、融了多少钱、投资方是谁
第2章 基础软件全栈——这是最核心的技术章节，华锐从操作系统到中间件到计算平台全栈自研，60纳秒的AMI消息中间件是所有产品的技术基座
第3章 核心交易平台ATP——覆盖110+券商的旗舰产品，从FPGA极速到全球交易的T1~T9完整产品线
第4章 实时风控ARC——30微秒穿透200条规则的流式内存计算引擎
第5章 高速行情AMD——FPGA硬件加速，130纳秒解码，行情速度决定交易速度
第6章 业务平台群——清算、投管、TA、交易总线等覆盖全链条的产品
第7章 信创与国际化——国产化替代和券商出海两大战略方向
第8章 总结与展望——核心结论和关键判断

每个章节都会深入到技术架构层面，让大家真正理解华锐的技术能力。我们开始吧。""")

# S3 Ch1 divider
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s,BG2)
bar(s,I(0),I(0),I(16),Pt(4),CY);bar(s,I(0),I(8.95),I(16),Pt(4),CY)
t(s,I(1.2),I(2.2),I(4),I(.6),"CHAPTER 01",18,CY,True)
bar(s,I(1.2),I(2.9),I(5),Pt(3),GO)
t(s,I(1.2),I(3.4),I(13),I(1.5),"公司概况与融资",44,W,True)
t(s,I(1.2),I(5),I(13),I(.6),"Company Overview & Funding",20,G);ftr(s)
n(s,"""【过渡页】
好，现在进入第一章，我们先来了解华锐技术的基本面——公司背景、团队规模、融资情况。这些信息帮助我们判断这家公司的实力和发展阶段。""")

# S4 Company Profile
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
hdr(s,"公司概况","深圳华锐分布式技术股份有限公司 · 2016年成立 · 500+员工")
card(s,I(.8),I(1.5),I(7.2),I(6.5))
mt(s,I(1.2),I(1.7),I(6.5),I(6),[
 ("公司全称",10,G),("深圳华锐分布式技术股份有限公司",15,W,True),("（2022年从'华锐金融技术'更名）",10,G),("",5,G),
 ("英文名",10,G),("Archforce Technology",15,CY,True),("",5,G),
 ("成立时间",10,G),("2016年",15,W,True),("",5,G),
 ("董事长 / CTO",10,G),("邹胜 / 赵楠",15,W,True),("",5,G),
 ("员工规模",10,G),("500+（含3位深圳市领军人才）",15,W,True),("",5,G),
 ("总部地址",10,G),("深圳市福田区新一代产业园5栋23层",14,W),
])
card(s,I(8.4),I(1.5),I(6.8),I(2.8))
t(s,I(8.7),I(1.65),I(6.2),I(.35),"使命 MISSION",12,CY,True)
t(s,I(8.7),I(2.05),I(6.2),I(.5),"提升中国金融IT基础设施自主能力",18,W,True)
bar(s,I(8.7),I(2.65),I(6),Pt(1),G)
t(s,I(8.7),I(2.8),I(6.2),I(.35),"愿景 VISION",12,CY,True)
t(s,I(8.7),I(3.2),I(6.2),I(.5),"分布式低时延技术引领者",18,W,True)
card(s,I(8.4),I(4.6),I(6.8),I(3.4))
t(s,I(8.7),I(4.75),I(6.2),I(.35),"五地布局",12,CY,True)
mt(s,I(8.7),I(5.2),I(6.2),I(2.5),[
 ("🏢 深圳（总部）— 福田区新一代产业园",12,W),
 ("🏢 上海 — 黄浦区世博滨江大厦",12,W),
 ("🏢 北京 — 朝阳区冠城大厦",12,W),
 ("🏢 长沙 — 岳麓区中建智慧谷",12,W),
 ("🏢 香港 — 大埔区香港科学园（2024年成立）",12,GO,True),
])
ftr(s)
n(s,"""【公司概况】
华锐技术的全称是"深圳华锐分布式技术股份有限公司"。注意这个名字——2022年从"华锐金融技术"更名为"华锐分布式技术"，把"分布式"写进了公司名，说明这是他们的核心技术定位。

公司2016年成立，到现在不到10年。董事长邹胜，CTO赵楠——赵楠曾在量化交易领域公开演讲，是技术领军人物。

500多名员工，其中包括3位深圳市领军人才。公司还聘请了柯伟祥先生作为资深战略顾问——他是中国资本市场登记结算业务的开创者。

使命是"提升中国金融IT基础设施自主能力"——这个定位非常清晰，就是做中国金融行业的核心技术底座。

特别值得注意的是，2024年在香港科学园成立了香港华锐，标志着国际化布局正式启动。目前在深圳、上海、北京、长沙、香港五个城市都有研发中心或子公司。""")

# S5 Funding
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
hdr(s,"融资情况","累计约 $83.19M（约6亿人民币）· D轮 · 18家投资方")
metric(s,I(.5),I(1.5),"$83M","累计融资总额",CY)
metric(s,I(3),I(1.5),"D轮","最新融资轮次",GO)
metric(s,I(5.5),I(1.5),"18家","投资方数量",GR)
metric(s,I(8),I(1.5),"2016","成立年份",W)
metric(s,I(10.5),I(1.5),"1.3亿","B轮（2020）",CY)
metric(s,I(13),I(1.5),"6亿+","人民币等值",GO)
bar(s,I(.8),I(3),I(14.4),Pt(2),CY)
t(s,I(.8),I(3.3),I(14),I(.4),"主要投资方 KEY INVESTORS",16,CY,True)
card(s,I(.8),I(3.8),I(4.5),I(2.2))
mt(s,I(1.1),I(4),I(4),I(1.8),[("🔵 红杉中国 Hongshan Capital",14,W,True),("顶级VC，成长性认可",11,G),("",4,G),("🔵 德诚基金 Decheng Fund",13,W),("🔵 亚布力投资 Yabuli",13,W)])
card(s,I(5.6),I(3.8),I(4.5),I(2.2))
mt(s,I(5.9),I(4),I(4),I(1.8),[("🟠 长江证券创新投资",14,W,True),("券商系资本，产业认可",11,G),("",4,G),("🟠 国信创投 Guoxin VC",13,W),("🟠 上海国信投资",13,W)])
card(s,I(10.4),I(3.8),I(4.8),I(2.2))
mt(s,I(10.7),I(4),I(4.2),I(1.8),[("⚪ Belay Capital",13,W),("⚪ Utrust Fund",13,W),("",4,G),("共18家投资方参与",13,G),("产业+财务双重认可",13,GO,True)])
card(s,I(.8),I(6.3),I(14.4),I(1.8))
mt(s,I(1.1),I(6.5),I(13.8),I(1.4),[
 ("💡 投资方特征分析",15,GO,True),
 ("B轮（2020年）：1.3亿人民币 → D轮（2022年12月）：金额未公开",13,W),
 ("券商系资本（长江证券、国信）= 行业内部高度认可华锐技术 + 顶级VC（红杉中国）= 高成长性认可",13,W),
 ("尚未上市，估值信息不透明。但多轮融资持续获得资本市场信心。",13,LG)])
ftr(s)
n(s,"""【融资情况】
根据CB Insights和PitchBook的数据，华锐技术累计融资约8319万美元，折合人民币约6亿元，已经完成了D轮融资。

投资方阵容非常值得关注，我把它分成三类来看：

第一类是顶级VC——红杉中国（现更名为Hongshan Capital）。红杉的参与说明华锐具有高成长性，符合顶级VC的投资标准。

第二类是券商系资本——长江证券创新投资、国信创投。这些是证券行业自己的资本，他们投华锐说明行业内部高度认可华锐的技术能力和产品价值。

第三类是其他财务投资人——德诚基金、亚布力投资、Belay Capital等，总共18家投资方。

已知的融资轮次：2020年B轮1.3亿人民币，2022年12月D轮金额未公开。

这个融资节奏说明华锐在快速成长，但还没有上市，估值信息不透明。不过从投资方阵容来看，产业资本+财务资本的双重认可，说明华锐的价值是被市场充分认可的。""")

# S6 Timeline
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
hdr(s,"发展历程 2016—2025","从创立到国际化，不到10年的快速成长")
events=[("2016","公司成立",CY),("2017","首个分布式实验室\nAMI/ATP/AMD V1.0",CY),("2018","上海华锐\n国家高新技术企业",W),("2019","北京华锐\n华锐研究所",W),("2020","B轮1.3亿\n长沙华锐",GO),("2021","国泰君安+华泰\n战略合作",W),("2022","小巨人 | D轮\nACP/AMQ/ATA/AIM",GO),("2023","T1/T3/T5/T7/D1\nIDC FinTech 50",CY),("2024","重点小巨人|香港\nT9/D7/D9|Forbes",GO),("2025","WAIC AI风控\n参编国家标准",GR)]
for i,(yr,ev,c) in enumerate(events):
 x=I(.3+i*1.54);card(s,x,I(1.5),I(1.44),I(6.5))
 t(s,x+I(.05),I(1.6),I(1.3),I(.4),yr,17,c,True,A.CENTER)
 bar(s,x+I(.2),I(2.05),I(1),Pt(2),c)
 mt(s,x+I(.08),I(2.2),I(1.28),I(4.5),[(l,9,W) for l in ev.split('\n')])
ftr(s)
n(s,"""【发展历程】
华锐的发展速度非常快，我挑几个关键节点讲：

2017年是奠基之年——创立了证券行业首个分布式技术实验室，这个实验室诞生在深圳车公庙的一栋旧厂房里。同年发布了AMI、ATP、AMD三大核心产品的V1.0版本。从旧厂房到行业标准，这个起点很有故事感。

2020-2021年是加速期——B轮融资1.3亿，成立长沙华锐。更重要的是2021年与国泰君安、华泰证券签订战略合作——这两家都是TOP10券商，说明华锐的产品已经获得了头部客户的认可。

2022年是转折年——更名为"深圳华锐分布式技术股份有限公司"，获评国家级专精特新小巨人，完成D轮融资，同时发布了ACP、AMQ、ATA、AIM四款新产品。产品线从核心交易扩展到了清算、投管、TA等全链条。

2023-2024年是爆发期——2023年发布T1/T3/T5/T7/D1/M9等多款新品，2024年升级为"重点小巨人"，成立香港华锐，发布T9全球交易平台，入选Forbes China Top 10。

2025年最新动态——在WAIC展示AI资管合规风控系统，作为证券行业唯一代表参编国家标准《消息中间件技术要求》。""")

prs.save('/Users/jiasunm/Code/archforce-research-v3/ppt/p1.pptx')
print(f"V3 P1: {len(prs.slides)} slides")
